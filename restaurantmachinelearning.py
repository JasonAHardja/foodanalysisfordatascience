import pandas as pd
import numpy as np
from datetime import datetime
from sklearn.cluster import KMeans
from sklearn.preprocessing import StandardScaler
import matplotlib.pyplot as plt
import os
from pptx import Presentation
from pptx.util import Inches
import shutil

def process_food_data(file_path, category):
    """
    Process the food data, clean it, and calculate profit.
    """
    df = pd.read_csv(file_path, engine='python', encoding='ISO-8859-1', sep=";", skiprows=2, on_bad_lines='skip')

    # Normalize column names: lowercase and strip whitespace
    df.columns = df.columns.str.strip().str.lower()

    # Define the expected column names
    menu_category_col = "menu category"
    menu_col = "menu"
    qty_col = "qty"
    total_col = "total"
    nett_sales_col = "nett sales"

    if menu_category_col not in df.columns or menu_col not in df.columns or qty_col not in df.columns or total_col not in df.columns or nett_sales_col not in df.columns:
        raise ValueError("Required columns missing.")

    df_filtered = df.loc[df[menu_category_col].str.upper() == category].copy()
    df_filtered["original_menu"] = df_filtered[menu_col]
    df_filtered[total_col] = pd.to_numeric(df_filtered[total_col].str.replace(".", "").str.replace(",", "."), errors='coerce').fillna(0)
    df_filtered[nett_sales_col] = pd.to_numeric(df_filtered[nett_sales_col].str.replace(".", "").str.replace(",", "."), errors='coerce').fillna(0)
    df_filtered["profit"] = df_filtered[total_col] - df_filtered[nett_sales_col]

    grouped_df = df_filtered.groupby("original_menu")[[qty_col, "profit"]].sum().reset_index()
    grouped_df = grouped_df.rename(columns={"original_menu": "menu", qty_col: "order amount"})
    return grouped_df, df_filtered

def perform_clustering(dataframe, n_clusters=3, category_name="Category", file_name="file"):
    """
    Perform K-Means clustering on 'order amount' and 'profit', and return the graph figure.
    """
    clustering_data = dataframe[['order amount', 'profit']].copy()
    scaler = StandardScaler()
    scaled_data = scaler.fit_transform(clustering_data)
    kmeans = KMeans(n_clusters=n_clusters, random_state=42)
    dataframe['Cluster'] = kmeans.fit_predict(scaled_data)

    # Create the graph
    plt.figure(figsize=(10, 6))
    plt.scatter(dataframe['order amount'], dataframe['profit'], c=dataframe['Cluster'], cmap='viridis', s=50)
    plt.xlabel("Order Amount")
    plt.ylabel("Profit")
    plt.title(f"K-Means Clustering of {category_name} Items\n({file_name})")
    plt.colorbar(label="Cluster")
    
    # Save figure and return it
    graph_file = f"{file_name}_{category_name}_clustering.png"
    plt.savefig(graph_file)
    plt.close()  # Close the plot to free up memory
    return graph_file

def save_to_csv(dataframe, output_file):
    dataframe.to_csv(output_file, index=False)
    print(f"Data saved to {output_file}")

def create_presentation(graph_files, output_pptx):
    # Create a PowerPoint presentation
    prs = Presentation()
    
    # Add a title slide (this is optional, you can remove this if not needed)
    slide_layout = prs.slide_layouts[0]  # Standard title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Clustering Analysis Results"
    
    # Add slides with clustering graphs and titles
    for graph_file, file_name in graph_files:
        slide_layout = prs.slide_layouts[6]  # Blank layout (no title)
        slide = prs.slides.add_slide(slide_layout)
        
        # Add a specific title for the clustering
        title_text = f"Clustering for {file_name}"
        
        # Create a text box for the title and place it at the top of the slide
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        
        # Add the graph image below the title
        slide.shapes.add_picture(graph_file, Inches(1), Inches(1.5), width=Inches(8.5), height=Inches(6))
    
    # Save the presentation
    prs.save(output_pptx)
    print(f"Presentation saved to {output_pptx}")



def process_multiple_files(file_paths, output_folder):
    """
    Process multiple CSV files and generate clustering results for each.
    """
    graph_files = []  # To store paths to saved graphs for presentation
    
    # Ensure the output folder exists
    os.makedirs(output_folder, exist_ok=True)
    
    for file_path in file_paths:
        try:
            # Process food and drink data
            foods_summary, foods_df = process_food_data(file_path, "MAKANAN")
            drinks_summary, drinks_df = process_food_data(file_path, "MINUMAN")

            # Get file name for the title
            file_name = os.path.basename(file_path).split('.')[0]
            
            # Perform clustering and get graph files
            clustered_foods_graph = perform_clustering(foods_summary, n_clusters=3, category_name="Food", file_name=file_name)
            clustered_drinks_graph = perform_clustering(drinks_summary, n_clusters=3, category_name="Drinks", file_name=file_name)
            
            graph_files.append((clustered_foods_graph, file_name))
            graph_files.append((clustered_drinks_graph, file_name))

            # Save CSVs to the output folder
            save_to_csv(foods_summary, os.path.join(output_folder, f"Clustered_Foods_{file_name}.csv"))
            save_to_csv(drinks_summary, os.path.join(output_folder, f"Clustered_Drinks_{file_name}.csv"))

            print(f"\nAnalysis complete for {file_name}! Generated files:")
            print(f"1. Clustered_Foods_{file_name}.csv")
            print(f"2. Clustered_Drinks_{file_name}.csv")
        
        except Exception as e:
            print(f"An error occurred while processing {file_path}: {str(e)}")
    
    # Create a PowerPoint presentation with the graphs
    pptx_file = os.path.join(output_folder, "Clustering_Results_Presentation.pptx")
    create_presentation(graph_files, pptx_file)
    print(f"PowerPoint presentation saved to {pptx_file}")
    
    # Move the output folder to the Desktop
    desktop_folder = os.path.join(os.path.expanduser("~"), "Desktop", "Clustering_Results")
    shutil.move(output_folder, desktop_folder)
    print(f"Output folder moved to Desktop: {desktop_folder}")

if __name__ == "__main__":
    file_paths = [
        "/Users/jasonhardjawidjaja/Desktop/4 Sales Recapitulation Detail Report april 2024.csv",
        "/Users/jasonhardjawidjaja/Desktop/3 Sales Recapitulation Detail Report maret 2024.csv",
        "/Users/jasonhardjawidjaja/Desktop/7 Sales Recapitulation Detail Report Juli 2024.csv",
        "/Users/jasonhardjawidjaja/Desktop/5 Sales Recapitulation Detail Report mei 2024.csv"
    ]
    
    # Create a folder on the Desktop to store the results
    output_folder = "/Users/jasonhardjawidjaja/Desktop/Clustering_Results"
    
    # Process multiple files and move the folder to Desktop
    process_multiple_files(file_paths, output_folder)
